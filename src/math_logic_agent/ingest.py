from __future__ import annotations

from multiprocessing import Process, Queue
from pathlib import Path
import time
from collections.abc import Callable

from pypdf import PdfReader
from pptx import Presentation

from .extraction import extract_equations, is_low_quality_text, ocr_pdf_page
from .models import RawDocument

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt", ".md"}


def discover_documents(input_dir: str | Path) -> list[Path]:
    root = Path(input_dir)
    if not root.exists():
        raise FileNotFoundError(f"Input directory does not exist: {root}")
    files: list[Path] = []
    for p in root.rglob("*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(p)
    return sorted(files)


def file_signature(path: str | Path) -> str:
    p = Path(path)
    stat = p.stat()
    return f"{stat.st_size}:{stat.st_mtime_ns}"


def ingest_path(
    path: str | Path,
    enable_ocr_fallback: bool = True,
    module_id: str | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> list[RawDocument]:
    p = Path(path)
    ext = p.suffix.lower()
    if ext == ".pdf":
        return ingest_pdf(
            p,
            enable_ocr_fallback=enable_ocr_fallback,
            module_id=module_id,
            progress_callback=progress_callback,
        )
    if ext == ".pptx":
        return ingest_pptx(p, module_id=module_id, progress_callback=progress_callback)
    if ext in {".txt", ".md"}:
        return ingest_text(p, module_id=module_id, progress_callback=progress_callback)
    return []


def ingest_path_safe(
    path: str | Path,
    enable_ocr_fallback: bool = True,
    module_id: str | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> tuple[list[RawDocument], str | None]:
    try:
        docs = ingest_path(
            path,
            enable_ocr_fallback=enable_ocr_fallback,
            module_id=module_id,
            progress_callback=progress_callback,
        )
        return docs, None
    except Exception as e:
        return [], f"{type(e).__name__}: {e}"


def _ingest_path_safe_worker(
    path: str,
    enable_ocr_fallback: bool,
    module_id: str | None,
    out_queue: Queue,
) -> None:
    def _heartbeat() -> None:
        out_queue.put(("heartbeat", None))

    docs, err = ingest_path_safe(
        path,
        enable_ocr_fallback=enable_ocr_fallback,
        module_id=module_id,
        progress_callback=_heartbeat,
    )
    out_queue.put(("result", (docs, err)))


def ingest_path_safe_with_timeout(
    path: str | Path,
    enable_ocr_fallback: bool = True,
    module_id: str | None = None,
    timeout_seconds: int = 0,
) -> tuple[list[RawDocument], str | None]:
    if timeout_seconds <= 0:
        return ingest_path_safe(path, enable_ocr_fallback=enable_ocr_fallback, module_id=module_id)

    out_queue: Queue = Queue(maxsize=1)
    proc = Process(
        target=_ingest_path_safe_worker,
        args=(str(path), enable_ocr_fallback, module_id, out_queue),
        daemon=True,
    )
    proc.start()
    last_progress = time.monotonic()

    while True:
        if not proc.is_alive() and out_queue.empty():
            return [], f"RuntimeError: ingest worker exited without result for {path}"

        if (time.monotonic() - last_progress) > timeout_seconds:
            proc.terminate()
            proc.join(2)
            return [], f"NoProgressTimeout: ingest made no progress for {timeout_seconds}s on {path}"

        try:
            kind, payload = out_queue.get(timeout=1)
        except Exception:
            continue

        if kind == "heartbeat":
            last_progress = time.monotonic()
            continue

        if kind == "result":
            proc.join(1)
            return payload

    return [], f"RuntimeError: unknown ingest worker state for {path}"


def ingest_pdf(
    path: Path,
    enable_ocr_fallback: bool = True,
    module_id: str | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> list[RawDocument]:
    try:
        reader = PdfReader(str(path), strict=False)
    except Exception as e:
        raise RuntimeError(f"Failed to open PDF: {path}") from e

    docs: list[RawDocument] = []
    page_errors = 0
    for i, page in enumerate(reader.pages, start=1):
        if progress_callback is not None:
            progress_callback()
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
            page_errors += 1
        text = text.strip()
        ocr_used = False
        if enable_ocr_fallback and is_low_quality_text(text):
            ocr_text = ocr_pdf_page(path, i)
            if ocr_text:
                text = ocr_text
                ocr_used = True
        if not text:
            continue
        equations = extract_equations(text)
        docs.append(
            RawDocument(
                text=text,
                source=str(path),
                module_id=module_id,
                page=i,
                metadata={
                    "filetype": "pdf",
                    "module_id": module_id,
                    "ocr_used": ocr_used,
                    "equation_count": len(equations),
                    "equations": equations,
                },
            )
        )
        if progress_callback is not None:
            progress_callback()

    if not docs and page_errors > 0:
        raise RuntimeError(f"All pages failed extraction for PDF: {path}")
    return docs


def ingest_pptx(
    path: Path,
    module_id: str | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> list[RawDocument]:
    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise RuntimeError(f"Failed to open PPTX: {path}") from e
    docs: list[RawDocument] = []
    for i, slide in enumerate(prs.slides, start=1):
        if progress_callback is not None:
            progress_callback()
        blocks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    blocks.append(t)
        if not blocks:
            continue
        docs.append(
            RawDocument(
                text="\n".join(blocks),
                source=str(path),
                module_id=module_id,
                page=i,
                metadata={"filetype": "pptx", "module_id": module_id},
            )
        )
        if progress_callback is not None:
            progress_callback()
    return docs


def ingest_text(
    path: Path,
    module_id: str | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> list[RawDocument]:
    if progress_callback is not None:
        progress_callback()
    try:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
    except Exception as e:
        raise RuntimeError(f"Failed to read text file: {path}") from e
    if not text:
        return []
    return [
        RawDocument(
            text=text,
            source=str(path),
            module_id=module_id,
            metadata={"filetype": path.suffix.lower().lstrip("."), "module_id": module_id},
        )
    ]


def ingest_directory(input_dir: str | Path, enable_ocr_fallback: bool = True) -> list[RawDocument]:
    all_docs: list[RawDocument] = []
    for path in discover_documents(input_dir):
        all_docs.extend(ingest_path(path, enable_ocr_fallback=enable_ocr_fallback))
    return all_docs


def ingest_paths(paths: list[Path], enable_ocr_fallback: bool = True, module_id: str | None = None) -> list[RawDocument]:
    all_docs: list[RawDocument] = []
    for path in paths:
        all_docs.extend(ingest_path(path, enable_ocr_fallback=enable_ocr_fallback, module_id=module_id))
    return all_docs
