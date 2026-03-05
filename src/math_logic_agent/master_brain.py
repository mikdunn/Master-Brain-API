from __future__ import annotations

import os
from collections import Counter
from dataclasses import dataclass
import math
import logging
from pathlib import Path
import re
import tomllib
from typing import Any, cast

from pypdf import PdfReader
from pptx import Presentation


_MIN_ALIASES_PER_MODULE = 10
_MAX_ALIASES_PER_MODULE = 20
_TARGET_ALIASES_PER_MODULE = 17
_MAX_FILES_PER_MODULE_FOR_TERMS = 0
_MAX_HEADINGS_PER_FILE = 40
_MAX_PDF_FILES_PER_MODULE_FOR_TERMS = 0
_MAX_PDF_SIZE_BYTES_FOR_TOC = 0
_PDF_FRONT_WORDS_LIMIT = 0
_PDF_FRONT_MAX_PAGES = 0

logging.getLogger("pypdf").setLevel(logging.ERROR)

_STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "by",
    "for",
    "from",
    "in",
    "into",
    "is",
    "it",
    "of",
    "on",
    "or",
    "that",
    "the",
    "their",
    "this",
    "to",
    "with",
    "day",
    "chapter",
    "section",
    "part",
    "notes",
    "study",
    "guide",
    "overview",
}

_NOISE_TOKENS = {
    "onlinelibrary",
    "wiley",
    "berklee",
    "stanford",
    "pearson",
    "cengage",
    "cambridge",
    "university",
    "press",
    "publisher",
    "publishers",
    "episode",
    "episodes",
    "notes",
    "textbook",
    "edition",
    "ed",
    "doi",
    "isbn",
    "copyright",
    "rights",
    "reserved",
}

_active_noise_tokens: set[str] | None = None

_GRAD_VOCAB_BY_BRAIN: dict[str, set[str]] = {
    "math_brain": {
        "variational",
        "convex",
        "optimization",
        "stability",
        "spectral",
        "eigen",
        "tensor",
        "topology",
        "proof",
    },
    "physics_brain": {
        "lagrangian",
        "hamiltonian",
        "quantum",
        "field",
        "renormalization",
        "symmetry",
        "gauge",
        "statistical",
        "dynamics",
    },
    "engineering_brain": {
        "control",
        "state-space",
        "estimation",
        "robust",
        "optimization",
        "dynamics",
        "signal",
        "system",
        "identification",
    },
    "science_brain": {
        "mechanism",
        "pathway",
        "regulation",
        "omics",
        "systems",
        "inference",
        "causal",
        "quantitative",
        "stochastic",
        "genomic",
        "proteomic",
        "epigenetic",
        "transcriptomic",
    },
    "business_brain": {
        "econometrics",
        "causal",
        "inference",
        "portfolio",
        "derivative",
        "volatility",
        "optimization",
        "forecast",
        "risk",
    },
    "cs_brain": {
        "algorithmic",
        "complexity",
        "distributed",
        "probabilistic",
        "bayesian",
        "inference",
        "transformer",
        "representation",
        "optimization",
    },
    "humanities_brain": {
        "historiography",
        "hermeneutic",
        "epistemology",
        "methodology",
        "rhetoric",
        "critical theory",
        "philology",
        "ethics",
        "interpretive",
    },
}


def get_default_master_brain_root() -> Path:
    env_root = (os.getenv("MASTER_BRAIN_ROOT") or "").strip()
    if env_root:
        return Path(env_root).expanduser()
    return Path.cwd() / "Master Brain"


DEFAULT_MASTER_BRAIN_ROOT = get_default_master_brain_root()

_MASTER_BRAIN_RELATIVE_DIRS: tuple[str, ...] = (
    "0_Math_Brain",
    "0_Math_Brain/Algebra",
    "0_Math_Brain/Calculus",
    "0_Math_Brain/Linear_Algebra",
    "0_Math_Brain/Differential_Equations",
    "0_Math_Brain/Probability_and_Stochastic_Processes",
    "0_Math_Brain/Optimization",
    "0_Math_Brain/Numerical_Methods",
    "0_Math_Brain/Logic_and_Proofs",
    "1_Physics_Brain",
    "1_Physics_Brain/Classical_Mechanics",
    "1_Physics_Brain/Thermodynamics_and_Statistical_Mechanics",
    "1_Physics_Brain/Electromagnetism",
    "1_Physics_Brain/Quantum_Physics",
    "1_Physics_Brain/Optics",
    "1_Physics_Brain/Fluid_Dynamics",
    "1_Physics_Brain/Continuum_Mechanics",
    "2_Engineering_Brain",
    "2_Engineering_Brain/Mechanical_Engineering",
    "2_Engineering_Brain/Mechanical_Engineering/Statics",
    "2_Engineering_Brain/Mechanical_Engineering/Dynamics",
    "2_Engineering_Brain/Mechanical_Engineering/Materials",
    "2_Engineering_Brain/Mechanical_Engineering/Thermofluids",
    "2_Engineering_Brain/Electrical_Engineering",
    "2_Engineering_Brain/Electrical_Engineering/Circuits",
    "2_Engineering_Brain/Electrical_Engineering/Signals_and_Systems",
    "2_Engineering_Brain/Electrical_Engineering/Control_Theory",
    "2_Engineering_Brain/Electrical_Engineering/Electromagnetics",
    "2_Engineering_Brain/Chemical_Engineering",
    "2_Engineering_Brain/Chemical_Engineering/Transport",
    "2_Engineering_Brain/Chemical_Engineering/Reaction_Engineering",
    "2_Engineering_Brain/Chemical_Engineering/Thermodynamics",
    "2_Engineering_Brain/Chemical_Engineering/Process_Modeling",
    "2_Engineering_Brain/Biomedical_Engineering",
    "2_Engineering_Brain/Biomedical_Engineering/Biomechanics",
    "2_Engineering_Brain/Biomedical_Engineering/Biomaterials",
    "2_Engineering_Brain/Biomedical_Engineering/Medical_Imaging",
    "2_Engineering_Brain/Biomedical_Engineering/Bioinstrumentation",
    "3_Science_Brain",
    "3_Science_Brain/Biology",
    "3_Science_Brain/Biology/Molecular_Biology",
    "3_Science_Brain/Biology/Genetics",
    "3_Science_Brain/Biology/Genomics",
    "3_Science_Brain/Biology/Bioinformatics",
    "3_Science_Brain/Biology/Systems_Biology",
    "3_Science_Brain/Biology/Biophysics",
    "3_Science_Brain/Biology/Mathematical_Biology",
    "3_Science_Brain/Biology/Ecology",
    "3_Science_Brain/Biology/Evolution",
    "3_Science_Brain/Biology/Microbiology",
    "3_Science_Brain/Biology/Physiology",
    "3_Science_Brain/Biology/Developmental_Biology",
    "3_Science_Brain/Chemistry",
    "3_Science_Brain/Chemistry/General_Chemistry",
    "3_Science_Brain/Chemistry/Organic_Chemistry",
    "3_Science_Brain/Chemistry/Physical_Chemistry",
    "3_Science_Brain/Chemistry/Analytical_Chemistry",
    "3_Science_Brain/Chemistry/Biochemistry",
    "3_Science_Brain/Chemistry/Quantum_Chemistry",
    "3_Science_Brain/Earth_and_Environmental_Sciences",
    "3_Science_Brain/Earth_and_Environmental_Sciences/Geology",
    "3_Science_Brain/Earth_and_Environmental_Sciences/Geophysics",
    "3_Science_Brain/Earth_and_Environmental_Sciences/Geochemistry",
    "3_Science_Brain/Earth_and_Environmental_Sciences/Environmental_Science",
    "3_Science_Brain/Earth_and_Environmental_Sciences/Climate_Science",
    "3_Science_Brain/Physics_of_Life",
    "3_Science_Brain/Physics_of_Life/Cellular_Biophysics",
    "3_Science_Brain/Physics_of_Life/Molecular_Biophysics",
    "3_Science_Brain/Physics_of_Life/Physical_Biology",
    "4_Business_Brain",
    "4_Business_Brain/Accounting",
    "4_Business_Brain/Finance",
    "4_Business_Brain/Finance/Corporate_Finance",
    "4_Business_Brain/Finance/Asset_Pricing",
    "4_Business_Brain/Finance/Derivatives",
    "4_Business_Brain/Finance/Portfolio_Theory",
    "4_Business_Brain/Economics",
    "4_Business_Brain/Economics/Microeconomics",
    "4_Business_Brain/Economics/Macroeconomics",
    "4_Business_Brain/Economics/Game_Theory",
    "4_Business_Brain/Economics/Industrial_Organization",
    "4_Business_Brain/Econometrics",
    "4_Business_Brain/Econometrics/Statistical_Inference",
    "4_Business_Brain/Econometrics/Time_Series",
    "4_Business_Brain/Econometrics/Causal_Inference",
    "4_Business_Brain/Fintech",
    "4_Business_Brain/Technical_Analysis",
    "5_Computer_Science_Brain",
    "5_Computer_Science_Brain/Algorithms",
    "5_Computer_Science_Brain/Data_Structures",
    "5_Computer_Science_Brain/Operating_Systems",
    "5_Computer_Science_Brain/Compilers",
    "5_Computer_Science_Brain/Distributed_Systems",
    "5_Computer_Science_Brain/Machine_Learning",
    "5_Computer_Science_Brain/Deep_Learning",
    "5_Computer_Science_Brain/Reinforcement_Learning",
    "5_Computer_Science_Brain/Probabilistic_Modeling",
    "5_Computer_Science_Brain/Scientific_Computing",
    "6_Humanities_Brain",
    "6_Humanities_Brain/Philosophy",
    "6_Humanities_Brain/History",
    "6_Humanities_Brain/Literature",
    "6_Humanities_Brain/Linguistics",
    "6_Humanities_Brain/Arts",
    "6_Humanities_Brain/Religion",
)


@dataclass(frozen=True, slots=True)
class ScaffoldSummary:
    root: Path
    total_directories: int
    created_directories: int
    existing_directories: int


def scaffold_master_brain_structure(
    root: str | Path = DEFAULT_MASTER_BRAIN_ROOT,
) -> ScaffoldSummary:
    root_path = Path(root)
    root_path.mkdir(parents=True, exist_ok=True)

    created = 0
    existing = 0
    for rel in _MASTER_BRAIN_RELATIVE_DIRS:
        p = root_path / rel
        if p.exists():
            existing += 1
        else:
            p.mkdir(parents=True, exist_ok=True)
            created += 1

    return ScaffoldSummary(
        root=root_path,
        total_directories=len(_MASTER_BRAIN_RELATIVE_DIRS),
        created_directories=created,
        existing_directories=existing,
    )


@dataclass(frozen=True, slots=True)
class _BrainDefinition:
    module_id: str
    display_name: str
    aliases: tuple[str, ...]
    candidate_roots: tuple[str, ...]


@dataclass(frozen=True, slots=True)
class _ModuleRow:
    module_id: str
    display_name: str
    paths: list[str]
    priority: int
    brain_module_id: str
    depth: int
    primary_terms: list[str]
    context_terms: list[str]


_BRAIN_DEFINITIONS: tuple[_BrainDefinition, ...] = (
    _BrainDefinition(
        module_id="math_brain",
        display_name="Math Brain",
        aliases=("math", "proof", "calculus", "linear algebra"),
        candidate_roots=("0_Math_Brain", "Math Brain"),
    ),
    _BrainDefinition(
        module_id="physics_brain",
        display_name="Physics Brain",
        aliases=(
            "physics",
            "mechanics",
            "electromagnetism",
            "quantum",
        ),
        candidate_roots=(
            "1_Physics_Brain",
            "Physics Brain",
            "Science Brain/Physics",
        ),
    ),
    _BrainDefinition(
        module_id="engineering_brain",
        display_name="Engineering Brain",
        aliases=("engineering", "design", "control", "thermofluids"),
        candidate_roots=("2_Engineering_Brain", "Engineering Brain"),
    ),
    _BrainDefinition(
        module_id="science_brain",
        display_name="Science Brain",
        aliases=(
            "biology",
            "chemistry",
            "geology",
            "environment",
            "life science",
        ),
        candidate_roots=("3_Science_Brain", "Science Brain"),
    ),
    _BrainDefinition(
        module_id="business_brain",
        display_name="Business Brain",
        aliases=(
            "business",
            "finance",
            "economics",
            "accounting",
            "econometrics",
        ),
        candidate_roots=("4_Business_Brain", "Business Brain"),
    ),
    _BrainDefinition(
        module_id="cs_brain",
        display_name="Computer Science Brain",
        aliases=(
            "computer science",
            "algorithms",
            "ml",
            "deep learning",
            "systems",
        ),
        candidate_roots=(
            "5_Computer_Science_Brain",
            "Computer Science Brain",
        ),
    ),
    _BrainDefinition(
        module_id="humanities_brain",
        display_name="Humanities Brain",
        aliases=(
            "humanities",
            "philosophy",
            "history",
            "literature",
            "linguistics",
            "arts",
            "religion",
        ),
        candidate_roots=("6_Humanities_Brain", "Humanities Brain"),
    ),
)


def _to_slug(text: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", text.strip().lower())
    slug = re.sub(r"_+", "_", slug).strip("_")
    return slug or "module"


def _normalize_phrase(text: str) -> str:
    cleaned = text.replace("_", " ")
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" -_/")
    return cleaned


def _is_course_code(term: str) -> bool:
    return bool(re.fullmatch(r"[a-z]{2,}\d{3,}[a-z]*", term.strip().lower()))


def _contains_course_code_token(term: str) -> bool:
    tokens = re.findall(r"[a-zA-Z0-9]+", term.lower())
    return any(_is_course_code(tok) for tok in tokens)


def _is_low_vowel_symbol_token(token: str) -> bool:
    if not token.isalpha():
        return False
    if len(token) < 4 or len(token) > 6:
        return False
    vowels = sum(ch in "aeiou" for ch in token.lower())
    return vowels <= 1


def is_noisy_lab_shorthand_term(term: str) -> bool:
    _unused: str = term
    _ = _unused
    return False


def _is_structural_noise_term(term: str) -> bool:
    lowered = term.lower().strip()
    if not lowered:
        return True
    if " / " in lowered:
        return True
    if lowered[0] in "([{" or lowered.endswith(")"):
        return True
    if re.search(r"\([^)]{0,80}\)", lowered):
        return True

    tokens = re.findall(r"[a-z0-9]+", lowered)
    if not tokens:
        return True

    if any(re.fullmatch(r"\d+", tok) for tok in tokens):
        return True

    if any(re.fullmatch(r"\d+(?:st|nd|rd|th)", tok) for tok in tokens):
        return True

    if any(re.fullmatch(r"\d{2,}[a-z]+", tok) for tok in tokens):
        return True

    if tokens[0] in {
        "figure",
        "fig",
        "table",
        "chapter",
        "section",
        "appendix",
        "copyright",
        "isbn",
        "doi",
        "stanford",
    }:
        return True

    if all(tok in _STOPWORDS for tok in tokens):
        return True

    active_noise = _active_noise_tokens or _NOISE_TOKENS
    if any(tok in active_noise for tok in tokens):
        return True

    if len(tokens) >= 5 and any(
        tok in active_noise
        for tok in tokens
    ):
        return True

    return False


def _is_noisy_for_brain(term: str, brain_module_id: str | None) -> bool:
    if _is_structural_noise_term(term):
        return True
    if brain_module_id != "science_brain":
        return False

    lowered = term.lower().strip()
    tokens = re.findall(r"[a-z0-9]+", lowered)
    if not tokens:
        return False

    has_symbolic = any(_is_low_vowel_symbol_token(tok) for tok in tokens)
    if not has_symbolic:
        return False

    vocab = _GRAD_VOCAB_BY_BRAIN.get("science_brain", set())
    has_grad = any(v in lowered for v in vocab)
    return not has_grad


def _iter_term_files(paths: list[Path]) -> list[Path]:
    files: list[Path] = []
    seen: set[Path] = set()
    allowed = {".pdf", ".pptx", ".md", ".txt"}
    for root in paths:
        if not root.exists() or not root.is_dir():
            continue
        for p in root.rglob("*"):
            if (
                _MAX_FILES_PER_MODULE_FOR_TERMS > 0
                and len(files) >= _MAX_FILES_PER_MODULE_FOR_TERMS
            ):
                return files
            if not p.is_file() or p.suffix.lower() not in allowed:
                continue
            if p in seen:
                continue
            seen.add(p)
            files.append(p)
    return files


def _extract_pdf_outline_terms(path: Path) -> list[str]:
    terms: list[str] = []
    try:
        if (
            _MAX_PDF_SIZE_BYTES_FOR_TOC > 0
            and path.stat().st_size > _MAX_PDF_SIZE_BYTES_FOR_TOC
        ):
            return terms
    except Exception:
        return terms
    try:
        reader = PdfReader(str(path), strict=False)
        outline = cast(Any, getattr(reader, "outline", None) or [])
    except KeyboardInterrupt:
        return terms
    except Exception:
        return terms

    def visit(node: object) -> None:
        if len(terms) >= _MAX_HEADINGS_PER_FILE:
            return
        if isinstance(node, list):
            for child in cast(list[Any], node):
                visit(child)
            return
        title = getattr(node, "title", None)
        if isinstance(title, str) and title.strip():
            terms.append(title.strip())

    visit(outline)
    return terms[:_MAX_HEADINGS_PER_FILE]


def _top_ngram_terms(words: list[str], limit: int = 24) -> list[str]:
    if not words:
        return []

    uni: Counter[str] = Counter()
    bi: Counter[str] = Counter()
    tri: Counter[str] = Counter()

    n = len(words)
    for i in range(n):
        w1 = words[i]
        if len(w1) >= 3 and w1 not in _STOPWORDS:
            uni[w1] += 1
        if i + 1 < n:
            w2 = words[i + 1]
            if (
                len(w1) >= 3
                and len(w2) >= 3
                and w1 not in _STOPWORDS
                and w2 not in _STOPWORDS
            ):
                bi[f"{w1} {w2}"] += 1
        if i + 2 < n:
            w2 = words[i + 1]
            w3 = words[i + 2]
            if (
                len(w1) >= 3
                and len(w2) >= 3
                and len(w3) >= 3
                and w1 not in _STOPWORDS
                and w2 not in _STOPWORDS
                and w3 not in _STOPWORDS
            ):
                tri[f"{w1} {w2} {w3}"] += 1

    scored: list[tuple[float, str]] = []
    for term, c in uni.items():
        scored.append((float(c), term))
    for term, c in bi.items():
        scored.append((float(c) * 1.4, term))
    for term, c in tri.items():
        scored.append((float(c) * 1.8, term))

    scored.sort(key=lambda x: (-x[0], x[1]))
    out: list[str] = []
    seen: set[str] = set()
    for _, term in scored:
        if term in seen:
            continue
        if _contains_course_code_token(term):
            continue
        if _is_structural_noise_term(term):
            continue
        seen.add(term)
        out.append(term)
        if len(out) >= limit:
            break
    return out


def _extract_pdf_front_terms(path: Path) -> list[str]:
    try:
        if (
            _MAX_PDF_SIZE_BYTES_FOR_TOC > 0
            and path.stat().st_size > _MAX_PDF_SIZE_BYTES_FOR_TOC
        ):
            return []
    except Exception:
        return []

    try:
        reader = PdfReader(str(path), strict=False)
    except KeyboardInterrupt:
        return []
    except Exception:
        return []

    page_texts: list[str] = []
    for idx, page in enumerate(reader.pages):
        if _PDF_FRONT_MAX_PAGES > 0 and idx >= _PDF_FRONT_MAX_PAGES:
            break
        try:
            text = page.extract_text() or ""
        except KeyboardInterrupt:
            return []
        except Exception:
            continue
        page_texts.append(text)

    if not page_texts:
        return []

    per_page_lines: list[list[str]] = []
    line_page_counts: Counter[str] = Counter()
    for text in page_texts:
        raw_lines = [
            re.sub(r"\s+", " ", line).strip().lower()
            for line in text.splitlines()
        ]
        lines = [line for line in raw_lines if line]
        per_page_lines.append(lines)
        line_page_counts.update(set(lines))

    page_count = len(page_texts)
    repeated_threshold = max(2, math.ceil(0.35 * page_count))
    repeated_lines = {
        line
        for line, count in line_page_counts.items()
        if count >= repeated_threshold and 2 <= len(line.split()) <= 14
    }

    top_line_counts: Counter[str] = Counter()
    bottom_line_counts: Counter[str] = Counter()
    for lines in per_page_lines:
        if not lines:
            continue
        top_line_counts.update(lines[:3])
        bottom_line_counts.update(lines[-3:])

    edge_threshold = max(2, math.ceil(0.2 * page_count))
    repeated_top = {
        line
        for line, count in top_line_counts.items()
        if count >= edge_threshold
    }
    repeated_bottom = {
        line
        for line, count in bottom_line_counts.items()
        if count >= edge_threshold
    }

    tokens: list[str] = []
    for lines in per_page_lines:
        page_filtered: list[str] = []
        n = len(lines)
        for idx, line in enumerate(lines):
            if line in repeated_lines:
                continue
            if idx < 3 and line in repeated_top:
                continue
            if idx >= max(0, n - 3) and line in repeated_bottom:
                continue
            page_filtered.append(line)
        text = "\n".join(page_filtered)
        words = re.findall(r"[a-zA-Z][a-zA-Z0-9\-']*", text.lower())
        if words:
            tokens.extend(words)
        if (
            _PDF_FRONT_WORDS_LIMIT > 0
            and len(tokens) >= _PDF_FRONT_WORDS_LIMIT
        ):
            break

    if _PDF_FRONT_WORDS_LIMIT > 0:
        tokens = tokens[:_PDF_FRONT_WORDS_LIMIT]

    return _top_ngram_terms(tokens)


def _extract_pptx_heading_terms(path: Path) -> list[str]:
    terms: list[str] = []
    try:
        prs = Presentation(str(path))
    except Exception:
        return terms

    for slide in prs.slides:
        if len(terms) >= _MAX_HEADINGS_PER_FILE:
            break
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            title_text = getattr(title_shape, "text", "")
        else:
            title_text = ""
        if isinstance(title_text, str) and title_text.strip():
            terms.append(title_text.strip())
            continue
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip() if isinstance(text, str) else ""
            if text:
                first_line = text.splitlines()[0].strip()
                if first_line:
                    terms.append(first_line)
                    break
    return terms[:_MAX_HEADINGS_PER_FILE]


def _extract_text_heading_terms(path: Path) -> list[str]:
    terms: list[str] = []
    try:
        lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    except Exception:
        return terms

    for line in lines[:300]:
        s = line.strip()
        if not s:
            continue
        if s.startswith("#"):
            terms.append(s.lstrip("#").strip())
            continue
        if re.match(r"^\d+(?:\.\d+)*\s+", s):
            terms.append(s)
            continue
        if len(s.split()) <= 10 and s == s.title():
            terms.append(s)
        if len(terms) >= _MAX_HEADINGS_PER_FILE:
            break
    return terms[:_MAX_HEADINGS_PER_FILE]


def _phrase_candidates(raw: str) -> list[str]:
    cleaned = _normalize_phrase(raw)
    cleaned = re.sub(r"[^a-zA-Z0-9\s\-/&,()]", " ", cleaned)
    cleaned = cleaned.replace("/", " ")
    cleaned = cleaned.replace("&", " and ")
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if not cleaned:
        return []

    words = [
        w.strip("()[]{}\"'`.,;:!?").lower()
        for w in cleaned.split()
    ]
    words = [w for w in words if len(w) >= 2]
    words = [w for w in words if w not in _STOPWORDS]
    if not words:
        return []

    out: list[str] = []
    n = len(words)
    for size in (1, 2, 3):
        if n < size:
            continue
        for i in range(0, n - size + 1):
            phrase_words = words[i:i + size]
            if all(w in _STOPWORDS for w in phrase_words):
                continue
            phrase = " ".join(phrase_words)
            if _is_course_code(phrase):
                continue
            if _is_structural_noise_term(phrase):
                continue
            out.append(phrase)
    return out


def _module_term_counter(paths: list[Path]) -> Counter[str]:
    counter: Counter[str] = Counter()
    pdf_seen = 0
    for file_path in _iter_term_files(paths):
        headings: list[str] = [file_path.stem]
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            if (
                _MAX_PDF_FILES_PER_MODULE_FOR_TERMS > 0
                and pdf_seen >= _MAX_PDF_FILES_PER_MODULE_FOR_TERMS
            ):
                continue
            pdf_seen += 1
            headings.extend(_extract_pdf_outline_terms(file_path))
            headings.extend(_extract_pdf_front_terms(file_path))
        elif suffix == ".pptx":
            headings.extend(_extract_pptx_heading_terms(file_path))
        else:
            headings.extend(_extract_text_heading_terms(file_path))

        for heading in headings[:_MAX_HEADINGS_PER_FILE]:
            for phrase in _phrase_candidates(heading):
                counter[phrase] += 1
    return counter


def select_specific_terms_for_aliases(
    module_counts: Counter[str],
    global_df: Counter[str],
    brain_module_id: str,
    depth: int,
    *,
    limit: int = 6,
) -> list[str]:
    scored: list[tuple[float, int, str]] = []
    vocab = _GRAD_VOCAB_BY_BRAIN.get(brain_module_id, set())
    depth_weight = 1.0 + (0.35 * max(depth, 0))

    for term, tf in module_counts.items():
        if len(term) < 4:
            continue
        if _is_course_code(term):
            continue
        if _is_noisy_for_brain(term, brain_module_id):
            continue
        df = global_df.get(term, 1)
        base = tf / float(df)
        lowered = term.lower()
        grad_hits = sum(1 for token in vocab if token in lowered)
        grad_bonus = float(grad_hits) * 1.15 * depth_weight
        score = base + grad_bonus
        scored.append((score, tf, term))
    scored.sort(key=lambda x: (-x[0], -x[1], x[2]))
    return [term for _, _, term in scored[:limit]]


def _build_aliases(
    primary_terms: list[str],
    context_terms: list[str],
    brain_module_id: str | None = None,
) -> list[str]:
    candidates: list[str] = []
    seen: set[str] = set()
    semantic_seen: set[str] = set()

    def canonical_key(term: str) -> str:
        tokens = re.findall(r"[a-z0-9]+", term.lower())
        active_noise = _active_noise_tokens or _NOISE_TOKENS
        filtered = [
            tok
            for tok in tokens
            if tok not in _STOPWORDS
            and tok not in active_noise
            and not re.fullmatch(r"\d+", tok)
            and not re.fullmatch(r"\d+(?:st|nd|rd|th)", tok)
            and not re.fullmatch(r"\d{2,}[a-z]+", tok)
        ]
        return " ".join(filtered)

    def add(term: str) -> None:
        normalized = _normalize_phrase(term).lower()
        if not normalized:
            return
        if normalized in seen:
            return
        if (
            _is_course_code(normalized)
            or _contains_course_code_token(normalized)
            or _is_structural_noise_term(normalized)
            or _is_noisy_for_brain(normalized, brain_module_id)
        ):
            return
        if len(normalized) > 120:
            return

        key = canonical_key(normalized)
        if not key:
            return
        if key in semantic_seen:
            return

        semantic_seen.add(key)
        seen.add(normalized)
        candidates.append(normalized)

    primaries = [
        _normalize_phrase(t)
        for t in primary_terms
        if _normalize_phrase(t)
    ]
    contexts = [
        _normalize_phrase(t)
        for t in context_terms
        if _normalize_phrase(t)
    ]

    for term in primaries:
        add(term)

    for term in primaries:
        add(f"{term} guide")
        add(f"{term} overview")
        add(f"{term} concepts")
        add(f"{term} fundamentals")
        add(f"{term} reference")
        add(f"learn {term}")
        add(f"study {term}")
        add(f"{term} practice")

    for ctx in contexts:
        add(ctx)
        for term in primaries:
            if ctx != term:
                add(f"{ctx} {term}")

    for term in primaries:
        add(f"master brain {term}")

    if candidates:
        pivot = candidates[0]
    else:
        pivot = "topic"

    pad = 1
    while len(candidates) < _MIN_ALIASES_PER_MODULE:
        add(f"{pivot} topic {pad}")
        pad += 1

    return candidates[:_TARGET_ALIASES_PER_MODULE]


def _toml_list(items: list[str], *, indent: int = 2) -> str:
    pad = " " * indent
    return "[\n" + "\n".join(f'{pad}"{item}",' for item in items) + "\n]"


def _normalize_noise_term(term: str) -> str:
    normalized = re.sub(r"\s+", " ", term.strip().lower())
    normalized = normalized.strip("\"'`.,;:!?()[]{}")
    return normalized


def _read_alias_noise_terms(config_path: Path) -> list[str]:
    if not config_path.exists():
        return []
    try:
        with config_path.open("rb") as f:
            data = tomllib.load(f)
    except Exception:
        return []

    raw_obj = data.get("alias_noise", {})
    raw = cast(dict[str, Any], raw_obj) if isinstance(raw_obj, dict) else {}
    terms_obj = raw.get("terms", [])
    terms = cast(list[Any], terms_obj) if isinstance(terms_obj, list) else []

    out: list[str] = []
    seen: set[str] = set()
    for item in terms:
        if not isinstance(item, str):
            continue
        term = _normalize_noise_term(item)
        if not term or term in seen:
            continue
        seen.add(term)
        out.append(term)
    return out


def _resolve_brain_roots(
    root: Path,
    candidates: tuple[str, ...],
) -> list[Path]:
    resolved = [root / rel for rel in candidates]
    existing = [p for p in resolved if p.exists() and p.is_dir()]
    return existing


def _collect_descendant_modules(
    roots: list[Path],
    top_module_id: str,
) -> list[tuple[str, str, list[str]]]:
    discovered: dict[tuple[str, ...], list[Path]] = {}
    labels: dict[tuple[str, ...], tuple[str, ...]] = {}

    for root in roots:
        if not root.exists() or not root.is_dir():
            continue
        for child in sorted(p for p in root.rglob("*") if p.is_dir()):
            rel_parts = child.relative_to(root).parts
            if not rel_parts:
                continue
            key = tuple(_to_slug(part) for part in rel_parts)
            discovered.setdefault(key, []).append(child)
            labels.setdefault(key, rel_parts)

    rows: list[tuple[str, str, list[str]]] = []
    for key in sorted(discovered.keys()):
        rel_labels = labels[key]
        module_id = f"{top_module_id}_{'_'.join(key)}"
        display_name = " / ".join(rel_labels)
        paths = [
            p.as_posix()
            for p in sorted(set(discovered[key]))
            if p.exists() and p.is_dir()
        ]
        rows.append((module_id, display_name, paths))
    return rows


def _render_module_block(
    module_id: str,
    display_name: str,
    paths: list[str],
    enabled: bool,
    stage: str,
    priority: int,
    aliases: list[str],
) -> str:
    lines = [
        f"[modules.{module_id}]",
        f'display_name = "{display_name}"',
        f"paths = {_toml_list(paths)}",
        f"enabled = {'true' if enabled else 'false'}",
        f'stage = "{stage}"',
        f"priority = {priority}",
        f"aliases = {_toml_list(aliases)}",
        "",
    ]
    return "\n".join(lines)


def render_master_module_registry_toml(
    root: str | Path = DEFAULT_MASTER_BRAIN_ROOT,
    custom_noise_terms: list[str] | None = None,
) -> str:
    global _active_noise_tokens

    base_noise = {_normalize_noise_term(t) for t in _NOISE_TOKENS}
    extra_noise = {
        _normalize_noise_term(t)
        for t in (custom_noise_terms or [])
        if _normalize_noise_term(t)
    }
    merged_noise = set(sorted(base_noise | extra_noise))

    prev_noise = _active_noise_tokens
    _active_noise_tokens = merged_noise

    r = Path(root)
    priority = 10
    chunks: list[str] = [
        "schema_version = 1",
        "",
        "[alias_noise]",
        f"terms = {_toml_list(sorted(merged_noise))}",
        "",
    ]
    rows: list[_ModuleRow] = []
    try:
        for brain in _BRAIN_DEFINITIONS:
            roots = _resolve_brain_roots(r, brain.candidate_roots)
            top_paths = [
                p.as_posix() for p in roots if p.exists() and p.is_dir()
            ]
            rows.append(
                _ModuleRow(
                    module_id=brain.module_id,
                    display_name=brain.display_name,
                    paths=top_paths,
                    priority=priority,
                    brain_module_id=brain.module_id,
                    depth=0,
                    primary_terms=[brain.display_name, *brain.aliases],
                    context_terms=[brain.module_id.replace("_", " ")],
                )
            )
            priority += 10

            descendant_modules = _collect_descendant_modules(
                roots,
                brain.module_id,
            )
            for module_id, display_name, paths in descendant_modules:
                local_depth = display_name.count("/") + 1
                alias_seed = [display_name, brain.display_name]
                alias_seed.extend(list(brain.aliases))
                alias_seed.extend(
                    part for part in display_name.split(" / ") if part
                )
                rows.append(
                    _ModuleRow(
                        module_id=module_id,
                        display_name=(
                            f"{brain.display_name} / {display_name}"
                        ),
                        paths=paths,
                        priority=priority,
                        brain_module_id=brain.module_id,
                        depth=local_depth,
                        primary_terms=alias_seed,
                        context_terms=[module_id.replace("_", " ")],
                    )
                )
                priority += 1

        module_term_counts: dict[str, Counter[str]] = {}
        global_df: Counter[str] = Counter()
        for row in rows:
            existing_paths = [Path(p) for p in row.paths if Path(p).exists()]
            counts = _module_term_counter(existing_paths)
            module_term_counts[row.module_id] = counts
            for term in counts.keys():
                global_df[term] += 1

        for row in rows:
            counts = module_term_counts.get(row.module_id, Counter())
            specific_terms = select_specific_terms_for_aliases(
                counts,
                global_df,
                brain_module_id=row.brain_module_id,
                depth=row.depth,
                limit=6,
            )
            generic_seed = row.primary_terms[:6]
            aliases = _build_aliases(
                primary_terms=[*specific_terms, *generic_seed],
                context_terms=row.context_terms,
                brain_module_id=row.brain_module_id,
            )
            chunks.append(
                _render_module_block(
                    module_id=row.module_id,
                    display_name=row.display_name,
                    paths=row.paths,
                    enabled=True,
                    stage="active",
                    priority=row.priority,
                    aliases=aliases,
                )
            )

        return "\n".join(chunks).rstrip() + "\n"
    finally:
        _active_noise_tokens = prev_noise


def write_master_module_registry(
    config_path: str | Path,
    root: str | Path = DEFAULT_MASTER_BRAIN_ROOT,
    overwrite: bool = True,
) -> Path:
    out = Path(config_path)
    if out.exists() and not overwrite:
        return out
    custom_noise_terms = _read_alias_noise_terms(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(
        render_master_module_registry_toml(
            root,
            custom_noise_terms=custom_noise_terms,
        ),
        encoding="utf-8",
    )
    return out
